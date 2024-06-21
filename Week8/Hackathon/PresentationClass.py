from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE


class PresentationClass:
    @staticmethod
    def create_presentation():
        # Create a presentation object
        prs = Presentation()

        # Slide 0: Title Slide
        slide_title = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide_title.shapes.title
        subtitle = slide_title.placeholders[1]
        title.text = ("Jacket Sales Report for 2024 (January-May)")
        subtitle.text = "Presentation on Sales Analysis"

        # Slide 1: Introduction
        slide_1 = prs.slides.add_slide(prs.slide_layouts[5])  # Choose layout 1 (Title and Content)
        slide_1.shapes.title.text = "Introduction"
        textbox = slide_1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1.5))
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        textbox.text = ("The object of the research is a company working with sales on\n "
                        "a well-known online platform in Russia (similar to Amazon).\n"
                        "The company sells clothes.\n"
                        "Four months of sales data was provided for a wide range of jackets, \n"
                        "including financial reports and customer feedbacks.\n"
                        "There was no task to make a financial report as this is a separate specialty. \n"
                        "Analyzed primarily qualitative rather than quantitative characteristics.\n")

        # Slide 2: Introduction
        slide_2 = prs.slides.add_slide(prs.slide_layouts[5])  # Choose layout 1 (Title and Content)
        slide_2.shapes.title.text = "Contents of the Columns:"
        textbox = slide_2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1.5))
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        textbox.text = ("Item Code: This column contains unique identifiers for items/products.\n"
                        "Brand: The brand name.\n"
                        "Supplier Article: The item (name).\n"
                        "Name: Name of the item.\n"
                        "Size: Size of the item purchased.\n"
                        "Payment Reason: The two values in this column are Sold and Returned.\n"
                        "Customer Order Date: Date when the customer placed the order.\n"
                        "Sales Date: Date when the item was sold.\n"
                        "Quantity: Number of items purchased/sold.\n"
                        "Retail Price: The highest price, more of a marketing thing to show bigger discounts \n"
                        "and customers would find items faster using the 'Discount' filter.\n"
                        "Wildberries Sold Goods (Pro): The price for which the online retailer sold the item.\n"
                        "Permanent Customer Discount (%): The customer's personal discount (from the store).\n"
                        "To Seller for Sold Goods: The price paid to the seller, after all duties and logistics fees have been collected.\n"
                        "Warehouse: Warehouse from which the item was dispatched or stored.\n"
                        "Country: Country associated with the transaction.\n"
                        "Barcode: Barcode associated with the item.\n"
                        "Size (Numerical): Numerical representation of the size.")

        # Define the slide layout
        title_and_content_layout = prs.slide_layouts[5]  # Layout 1 (Title and Content)

        # Dictionary containing titles and filenames for each plot
        titles_to_filenames = {
            "Total Sales Quantity by Brand": "plot_sales_quantity_by_brand.png",
            "Quantity and Percentage of Sales and Returns of Jackets for January-April 2024": "sales_returns_plot.png",
            "Sales for Brand": "plot_sales_for_brand.png",
            "Daily Order Dynamic for 4 months": "plot_daily_orders.png",
            "Scatterplot of Sales and Returns": "plot_scatterplot_sales_and_returns.png",
            "Distribution of Jacket Prices for Purchased and Returned Jackets": "plot_price_distribution.png",
            "Most Popular Jackets": "plot_most_popular_jackets.png",
            "Top 1": "image_138646319.jpg",
            "Top 2": "image_138646324.jpg",
            "Top 3": "image_146321949.jpg",
            "Top 4": "image_138646323.jpg",
            "Analyze Sales Performance": "plot_analyze_sales_performance.png",
            "Scatterplot Sales vs. Returns for Each Jacket for Each Size": "scatterplot_sales_vs_returns.png",
            "Total Sales Quantity by Size": "plot_sales_quantity_by_size.png",
            "Subplot of Top 4 Sales vs Returns": "subplot_top4_sales_vs_returns.png",
            "Monthly Size of Transactions": "plot_monthly_size.png",
            "Boxplots of Top 4 Sales": "plot_boxplots_top4_sales.png",
            "Analyze Return Percentage": "plot_analyze_return_percentage.png",
            "Histogram Plots of Top 4 Sales": "plot_histplots_top4_sales.png",
            "Quantity per Day of Week": "plot_quantity_per_day_of_week.png",
            "Frequency of Sales and Average Discount by Discount Lavel": "categorize_discount_1.png",
            "Return Rate by Discount Level": "categorize_discount_2.png",
            "Top 10 Warehouses by Jacket Sales": "plot_warehouses_1.png",
            "Top 10 Warehouses by Jacket Retuns": "plot_warehouses_2.png",
            "Total Deliveries and Returns Percentage by Warehouse": "plot_warehouses_3.png",
            "Correlation of Returns": "plot_correlation_returns.png",
            "Distribution of Star Rating": "plot_star_count.png",
            "Correlation Heatmap": "create_common_df.png",
            "Distribution of Sentiment Categories (SentimentIntensityAnalyzer)": "sentiment_sid_pie_chart.png",
            "Distribution of Sentiment Categories (TextBlob)": "sentiment_textblob_pie_chart.png",
            "Distribution of Sentiments (My Version)": "sentiment_my_pie_chart.png",
            "Word Cloud of Feedback Text": "cloud_feedbacks_1.png",
            "Word Cloud of Feedback Text (Without Russian Stopwords)": "cloud_feedbacks_2.png",
            "Customized Scatter Plot for Top 20 Most Common Words in Feedbacks with Different Star Counts": "scatter_plot_top_20.png",
            "Venn Diagram of Top Keywords": "venn_diagram.png"

        }
        # Loop through the dictionary items and add slides for each plot
        for title, filename in titles_to_filenames.items():           # Add a slide for the current function
            slide = prs.slides.add_slide(title_and_content_layout)

            # Set the slide title
            slide.shapes.title.text = title

            # Add the picture (plot) to the slide
            slide.shapes.add_picture(filename, Inches(1), Inches(2), width=Inches(8))


        # # Slide 7: Daily Orders Plot
        # slide_7 = prs.slides.add_slide(prs.slide_layouts[5])  # Choose layout 1 (Title and Content)
        # slide_7.shapes.title.text = "Distribution of Jacket Prices for Purchased and Retuned Jackets"
        # slide_7.shapes.add_picture('plot_price_distribution.png', Inches(1), Inches(2), width=Inches(8))

        # Slide N: Introduction
        slide_N = prs.slides.add_slide(prs.slide_layouts[5])  # Choose layout 1 (Title and Content)
        slide_N.shapes.title.text = "Customer profile:"
        textbox = slide_N.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1.5))
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        textbox.text = ("A woman from Russia, Belarus or Kazakhstan named Elena, Olga or Natalia \n"
                        "with a personal discount percentage of 15-20%, size 46, 48 or 44, \n"
                        "from the whole range of the oversize blazer she prefers \n"
                        "brown, green, gray, black colors. \n"
                        "She buys less often on Thursdays and Mondays than on Wednesdays and Sundays. \n"
                        "She most often leaves positive feedback, \n"
                        "although she thinks that the blazers are too big and it is necessary \n"
                        "to pay attention to the size grid, material and she is concerned about the sleeves."
                        "Common words among all star ratings: \т"
                        "{'fabric', 'sleeves', 'height', 'cm', 'size', 'fit', 'shape', 'larger'}")
        # Save the presentation
        prs.save("sales_analysis_presentation.pptx")